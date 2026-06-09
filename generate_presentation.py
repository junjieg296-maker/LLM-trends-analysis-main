# -*- coding: utf-8 -*-
"""
LLM Multi-Agent Research Trends Analysis - Professional Presentation
专业学术汇报PPT - 简洁优雅版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 设置路径
PROJECT_ROOT = r"F:\文献计量学\LLM-trends-analysis-main"
OUTPUT_DIR = os.path.join(PROJECT_ROOT, "outputs")
FIGURES_DIR = os.path.join(OUTPUT_DIR, "figures")
REVIEW_FIGURES_DIR = os.path.join(OUTPUT_DIR, "review_figures")
OUTPUT_PPT = os.path.join(PROJECT_ROOT, "LLM_Research_Presentation_Simple.pptx")

# PPT尺寸设置 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 简洁优雅配色方案 - 只用深蓝+白色+灰色
COLORS = {
    'primary': RgbColor(30, 58, 138),       # 深蓝 - 主色调
    'accent': RgbColor(30, 58, 138),        # 深蓝 - 强调色（统一）
    'dark': RgbColor(31, 31, 31),           # 深黑 - 主要文字
    'medium': RgbColor(107, 114, 128),      # 中灰 - 次要文字
    'light': RgbColor(243, 244, 246),       # 浅灰 - 背景
    'white': RgbColor(255, 255, 255),       # 白色
    'border': RgbColor(209, 213, 219),      # 边框灰
}

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 添加所有幻灯片
    add_title_slide(prs)
    add_overview_slide(prs)
    add_methodology_slide(prs)
    add_data_quality_slide(prs)
    add_rq1_slide(prs)
    add_rq2_slide(prs)
    add_rq3_slide(prs)
    add_rq4_slide(prs)
    add_rq5_slide(prs)
    add_conclusion_slide(prs)
    add_acknowledgment_slide(prs)

    # 保存
    prs.save(OUTPUT_PPT)
    print(f"PPT已生成: {OUTPUT_PPT}")
    return OUTPUT_PPT

def add_title_slide(prs):
    """封面页 - 简洁深蓝风格"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 深蓝背景
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']

    # 主标题
    main_title = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.5), Inches(1.8))
    tf = main_title.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "基于文献计量学的"
    p.font.size = Pt(26)
    p.font.color.rgb = RgbColor(180, 190, 210)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "LLM多智能体协作研究趋势分析"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 英文副标题
    sub_title = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.5), Inches(0.6))
    tf = sub_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Bibliometric Analysis of LLM Multi-Agent Collaboration Research"
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    # 分隔线
    sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(5), Inches(3.5), Pt(2))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = COLORS['white']
    sep_line.line.fill.background()

    # 项目信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.5), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "汇报人：高俊杰"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "课程：文献计量学、前沿趋势追踪与项目制学习框架"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(180, 190, 210)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "数据来源：Lens.org | 有效文献：522篇 | 时间跨度：2020-2026"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

def add_overview_slide(prs):
    """项目概览与RQ总览"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "研究问题与项目概览")

    # 左侧 - 数据集概览卡片
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.4), Inches(5.5), Inches(5.7))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLORS['light']
    left_card.line.color.rgb = COLORS['border']

    # 数据集标题
    data_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5), Inches(0.5))
    tf = data_title.text_frame
    p = tf.paragraphs[0]
    p.text = "数据集核心指标"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 数据指标 - 统一深蓝色数字
    metrics = [
        ("522", "有效文献", "清洗去重后纳入分析"),
        ("4,265", "总被引次数", "领域整体学术关注度"),
        ("31", "h指数", "有31篇论文被引≥31次"),
        ("8.17", "篇均被引", "高于CS领域平均水平"),
        ("181.29%", "年均增长率", "高速扩张期特征"),
    ]

    y = Inches(2.1)
    for num, label, desc in metrics:
        # 数字 - 统一深蓝
        num_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(1.6), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']

        # 标签 - 深黑
        label_box = slide.shapes.add_textbox(Inches(2.5), y, Inches(2.5), Inches(0.35))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(Inches(2.5), y + Inches(0.35), Inches(2.8), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.85)

    # 检索策略信息
    search_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(5.2), Inches(0.6))
    tf = search_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "检索策略：(Object) AND (Method) 布尔逻辑"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['medium']

    # 右侧 - 五个RQ详细展示
    rq_title = slide.shapes.add_textbox(Inches(6.2), Inches(1.5), Inches(6.5), Inches(0.5))
    tf = rq_title.text_frame
    p = tf.paragraphs[0]
    p.text = "研究问题 (Research Questions)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    rqs = [
        ("RQ1", "发文量变化与发展阶段", "分析2020-2026年发文量趋势，判断领域所处发展阶段"),
        ("RQ2", "核心研究主体与合作格局", "识别核心作者和机构，分析合作网络结构特征"),
        ("RQ3", "研究热点与技术结构", "揭示高频主题，分析研究主题之间的技术层次关系"),
        ("RQ4", "知识基础与代表性论文", "识别高被引论文，构建综述的技术叙事线"),
        ("RQ5", "研究空白与未来挑战", "基于计量证据，归纳领域未来研究方向"),
    ]

    y = Inches(2.1)
    for rq, title, desc in rqs:
        # RQ标签 - 统一深蓝
        rq_tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), y, Inches(0.85), Inches(0.45))
        rq_tag.fill.solid()
        rq_tag.fill.fore_color.rgb = COLORS['primary']
        rq_tag.line.fill.background()

        rq_text = slide.shapes.add_textbox(Inches(6.25), y + Inches(0.05), Inches(0.75), Inches(0.35))
        tf = rq_text.text_frame
        p = tf.paragraphs[0]
        p.text = rq
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 标题 - 深黑
        title_box = slide.shapes.add_textbox(Inches(7.15), y, Inches(5.5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(Inches(7.15), y + Inches(0.4), Inches(5.5), Inches(0.45))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']

        y += Inches(1)

def add_methodology_slide(prs):
    """方法论与研究框架"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "研究方法与框架")

    # 方法流程图
    img_path = os.path.join(REVIEW_FIGURES_DIR, "fig1_methodology_workflow.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.5), width=Inches(7.5))

    # 右侧 - 方法详解
    right_x = Inches(8.2)

    # 分析方法标题
    method_title = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(4.8), Inches(0.4))
    tf = method_title.text_frame
    p = tf.paragraphs[0]
    p.text = "分析方法体系"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 方法列表
    methods = [
        ("文献计量学", "年度趋势、影响力指标、h指数计算"),
        ("社会网络分析", "作者合作网络、关键词共现网络"),
        ("引用分析", "高被引论文识别、共被引网络构建"),
        ("参数敏感性测试", "阈值扫描、时间窗口比较"),
    ]

    y = Inches(2)
    for method, desc in methods:
        # 方法名称 - 深黑
        method_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.35))
        tf = method_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {method}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(right_x + Inches(0.3), y + Inches(0.35), Inches(4.5), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.7)

    # 工具栈
    tool_title = slide.shapes.add_textbox(right_x, Inches(4.8), Inches(4.8), Inches(0.4))
    tf = tool_title.text_frame
    p = tf.paragraphs[0]
    p.text = "技术工具栈"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    tools = [
        ("Python 3.8+", "数据处理与分析"),
        ("Pandas", "数据清洗与去重"),
        ("NetworkX", "网络构建与可视化"),
        ("Matplotlib", "图表生成"),
    ]

    y = Inches(5.25)
    for tool, use in tools:
        tool_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.35))
        tf = tool_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {tool} - {use}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['medium']
        y += Inches(0.35)

def add_data_quality_slide(prs):
    """数据质量与PRISMA流程"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "数据筛选与质量控制")

    # PRISMA流程图 - 统一深蓝
    # 识别阶段
    id_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3), Inches(1.8))
    id_box.fill.solid()
    id_box.fill.fore_color.rgb = COLORS['primary']
    id_box.line.fill.background()

    id_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(2.6), Inches(0.4))
    tf = id_title.text_frame
    p = tf.paragraphs[0]
    p.text = "识别 Identification"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    id_num = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(2.6), Inches(0.5))
    tf = id_num.text_frame
    p = tf.paragraphs[0]
    p.text = "542"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    id_desc = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(2.6), Inches(0.5))
    tf = id_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "Lens.org初始检索导出"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(200, 210, 230)

    # 筛选阶段
    screen_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(1.5), Inches(3), Inches(1.8))
    screen_box.fill.solid()
    screen_box.fill.fore_color.rgb = COLORS['primary']
    screen_box.line.fill.background()

    screen_title = slide.shapes.add_textbox(Inches(4.2), Inches(1.6), Inches(2.6), Inches(0.4))
    tf = screen_title.text_frame
    p = tf.paragraphs[0]
    p.text = "筛选 Screening"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    screen_num = slide.shapes.add_textbox(Inches(4.2), Inches(2.1), Inches(2.6), Inches(0.5))
    tf = screen_num.text_frame
    p = tf.paragraphs[0]
    p.text = "-20"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    screen_desc = slide.shapes.add_textbox(Inches(4.2), Inches(2.6), Inches(2.6), Inches(0.5))
    tf = screen_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "DOI+Title双重去重(3.69%)"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(200, 210, 230)

    # 纳入阶段
    include_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.5), Inches(3), Inches(1.8))
    include_box.fill.solid()
    include_box.fill.fore_color.rgb = COLORS['primary']
    include_box.line.fill.background()

    include_title = slide.shapes.add_textbox(Inches(7.7), Inches(1.6), Inches(2.6), Inches(0.4))
    tf = include_title.text_frame
    p = tf.paragraphs[0]
    p.text = "纳入 Included"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    include_num = slide.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(2.6), Inches(0.5))
    tf = include_num.text_frame
    p = tf.paragraphs[0]
    p.text = "522"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    include_desc = slide.shapes.add_textbox(Inches(7.7), Inches(2.6), Inches(2.6), Inches(0.5))
    tf = include_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "有效文献纳入分析"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(200, 210, 230)

    # 箭头
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(2.2), Inches(0.5), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLORS['medium']
    arrow1.line.fill.background()

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(2.2), Inches(0.5), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS['medium']
    arrow2.line.fill.background()

    # 字段覆盖率
    quality_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.5))
    tf = quality_title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心字段覆盖率"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 字段覆盖率表格
    fields = [
        ("Title (标题)", "100.00%", "完全满足词频与共现分析"),
        ("Year (年份)", "100.00%", "完全满足时序演进分析"),
        ("Abstract (摘要)", "98.28%", "满足深度语义挖掘与主题聚类"),
        ("References", "81.99%", "满足文献共被引网络构建"),
        ("DOI", "99.81%", "支持去重、溯源和milestone识别"),
    ]

    y = Inches(4.3)
    for field, rate, use in fields:
        # 字段名称 - 深黑
        field_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(3), Inches(0.4))
        tf = field_box.text_frame
        p = tf.paragraphs[0]
        p.text = field
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 覆盖率 - 深蓝
        rate_box = slide.shapes.add_textbox(Inches(3.5), y, Inches(1.5), Inches(0.4))
        tf = rate_box.text_frame
        p = tf.paragraphs[0]
        p.text = rate
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']

        # 用途 - 中灰
        use_box = slide.shapes.add_textbox(Inches(5), y, Inches(7.5), Inches(0.4))
        tf = use_box.text_frame
        p = tf.paragraphs[0]
        p.text = use
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.45)

def add_rq1_slide(prs):
    """RQ1：发文量变化与发展阶段"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "RQ1：发文量变化与发展阶段")

    # 左侧图片
    img_path = os.path.join(REVIEW_FIGURES_DIR, "fig2_publication_growth_stages.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.5), width=Inches(7.5))

    # 右侧详细分析
    right_x = Inches(8.2)

    # 阶段分析标题
    phase_title = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(4.8), Inches(0.4))
    tf = phase_title.text_frame
    p = tf.paragraphs[0]
    p.text = "发展阶段判断"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 阶段说明
    phases = [
        ("2020-2023", "萌芽期", "文献数量较少(1-19篇)，处于技术探索阶段，研究者开始关注LLM与多智能体结合的可能性"),
        ("2024", "起步期", "发文量上升至81篇，agentic AI、工具调用、RAG成为独立研究热点，领域开始形成"),
        ("2025", "爆发期", "发文量激增至344篇(同比增长325%)，领域进入高速扩张，大量新研究者涌入"),
    ]

    y = Inches(2)
    for year, phase, desc in phases:
        # 年份标签 - 深蓝
        year_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y, Inches(1.3), Inches(0.4))
        year_box.fill.solid()
        year_box.fill.fore_color.rgb = COLORS['primary']
        year_box.line.fill.background()

        year_text = slide.shapes.add_textbox(right_x + Inches(0.1), y + Inches(0.05), Inches(1.1), Inches(0.3))
        tf = year_text.text_frame
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 阶段名称 - 深黑
        phase_box = slide.shapes.add_textbox(right_x + Inches(1.4), y, Inches(1.2), Inches(0.4))
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(right_x, y + Inches(0.45), Inches(4.8), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['medium']

        y += Inches(1.1)

    # 关键指标
    indicator_title = slide.shapes.add_textbox(right_x, Inches(5.5), Inches(4.8), Inches(0.4))
    tf = indicator_title.text_frame
    p = tf.paragraphs[0]
    p.text = "关键影响力指标"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    indicators = [
        ("181.29%", "平均年增长率"),
        ("325%", "2025年同比增长"),
        ("4,265", "总被引次数"),
        ("31", "h指数"),
    ]

    y = Inches(5.95)
    for num, label in indicators:
        # 数字 - 深蓝
        num_box = slide.shapes.add_textbox(right_x, y, Inches(1.5), Inches(0.35))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']

        # 标签 - 中灰
        label_box = slide.shapes.add_textbox(right_x + Inches(1.6), y + Inches(0.05), Inches(3), Inches(0.25))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.35)

def add_rq2_slide(prs):
    """RQ2：核心研究主体与合作格局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "RQ2：核心研究主体与合作格局")

    # 作者合作网络图片
    img_path = os.path.join(FIGURES_DIR, "author_collaboration_network.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.5), width=Inches(7.5))

    # 右侧详细分析
    right_x = Inches(8.2)

    # 核心发现
    find_title = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(4.8), Inches(0.4))
    tf = find_title.text_frame
    p = tf.paragraphs[0]
    p.text = "核心发现"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 发现内容
    finds = [
        "研究力量高度分散，尚未形成稳定核心团队",
        "最高产作者仅发文3篇，说明领域仍处于快速扩张期",
        "跨学科特征明显，覆盖AI、机器人、人机交互、医学等",
        "合作以项目型为主，缺乏长期稳定的合作关系",
    ]

    y = Inches(2)
    for find in finds:
        find_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.45))
        tf = find_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {find}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']
        y += Inches(0.45)

    # 作者统计
    stats_title = slide.shapes.add_textbox(right_x, Inches(4), Inches(4.8), Inches(0.4))
    tf = stats_title.text_frame
    p = tf.paragraphs[0]
    p.text = "作者与合作统计"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    stats = [
        ("约 2,565", "作者总数"),
        ("约 5.1", "篇均作者数"),
        ("3 篇", "最高产作者发文量"),
        ("51 个", "过滤后合作节点"),
        ("48 条", "过滤后合作边"),
    ]

    y = Inches(4.45)
    for num, label in stats:
        # 数字 - 深蓝
        num_box = slide.shapes.add_textbox(right_x, y, Inches(1.8), Inches(0.35))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']

        # 标签 - 中灰
        label_box = slide.shapes.add_textbox(right_x + Inches(1.9), y + Inches(0.05), Inches(3), Inches(0.25))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.38)

    # 机构分布
    inst_title = slide.shapes.add_textbox(right_x, Inches(6.5), Inches(4.8), Inches(0.4))
    tf = inst_title.text_frame
    p = tf.paragraphs[0]
    p.text = "机构/国家分布"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    inst_desc = slide.shapes.add_textbox(right_x, Inches(6.9), Inches(4.8), Inches(0.4))
    tf = inst_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "北美、欧洲、亚洲均有研究力量参与，呈现全球化特征"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['medium']

def add_rq3_slide(prs):
    """RQ3：研究热点与技术结构"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "RQ3：研究热点与技术结构")

    # 关键词共现网络图片
    img_path = os.path.join(FIGURES_DIR, "keyword_cooccurrence_network.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.5), width=Inches(7.5))

    # 右侧详细分析
    right_x = Inches(8.2)

    # 四层技术体系
    layer_title = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(4.8), Inches(0.4))
    tf = layer_title.text_frame
    p = tf.paragraphs[0]
    p.text = "四层技术体系"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    layers = [
        ("1. 技术底座", "LLM, AI, Machine Learning, NLP", "多智能体系统依赖LLM和AI基础能力"),
        ("2. 协作架构", "Agentic AI, Multi-Agent Systems", "研究重点转向角色分工、任务拆解"),
        ("3. 能力增强", "RAG, RL, Planning, Tool Use", "关注知识检索、工具调用、复杂规划"),
        ("4. 治理应用", "Ethics, HRI, Digital Health, Robot", "进入高约束应用和安全治理场景"),
    ]

    y = Inches(2)
    for layer, keywords, desc in layers:
        # 层级名称 - 深黑
        name_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.35))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = layer
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 关键词 - 中灰
        kw_box = slide.shapes.add_textbox(right_x + Inches(0.3), y + Inches(0.35), Inches(4.5), Inches(0.35))
        tf = kw_box.text_frame
        p = tf.paragraphs[0]
        p.text = keywords
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.75)

    # 高频关键词
    kw_title = slide.shapes.add_textbox(right_x, Inches(5), Inches(4.8), Inches(0.4))
    tf = kw_title.text_frame
    p = tf.paragraphs[0]
    p.text = "高频关键词 TOP5"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    top_kw = [
        ("artificial intelligence", "56次"),
        ("large language models", "45次"),
        ("large language model", "24次"),
        ("computer science", "22次"),
        ("machine learning", "18次"),
    ]

    y = Inches(5.45)
    for kw, count in top_kw:
        kw_box = slide.shapes.add_textbox(right_x, y, Inches(3.5), Inches(0.3))
        tf = kw_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {kw}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']

        count_box = slide.shapes.add_textbox(Inches(11.2), y, Inches(1.2), Inches(0.3))
        tf = count_box.text_frame
        p = tf.paragraphs[0]
        p.text = count
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']

        y += Inches(0.3)

def add_rq4_slide(prs):
    """RQ4：知识基础与代表性论文"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "RQ4：知识基础与代表性论文")

    # milestone候选矩阵图片
    img_path = os.path.join(REVIEW_FIGURES_DIR, "fig4_milestone_roadmap.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.5), width=Inches(7.5))

    # 右侧详细分析
    right_x = Inches(8.2)

    # Top论文标题
    paper_title = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(4.8), Inches(0.4))
    tf = paper_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Top 5 高被引论文"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    papers = [
        ("[2024] 352次", "自驱动实验室、化学与材料科学", "AI驱动的科学发现闭环"),
        ("[2023] 259次", "无人机与多智能体系统综述", "多智能体协调机制"),
        ("[2023] 173次", "人机协同创作与语言模型", "人类监督与AI协作"),
        ("[2025] 136次", "化学领域LLM与autonomous agents", "垂直领域智能体应用"),
        ("[2024] 109次", "SMART-LLM多智能体任务规划", "自然语言任务分解"),
    ]

    y = Inches(2)
    for year_cite, topic, theme in papers:
        paper_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.55))
        tf = paper_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"• {year_cite} - {topic}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        p = tf.add_paragraph()
        p.text = f"   技术主线：{theme}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.6)

    # 四条技术主线
    theme_title = slide.shapes.add_textbox(right_x, Inches(5.2), Inches(4.8), Inches(0.4))
    tf = theme_title.text_frame
    p = tf.paragraphs[0]
    p.text = "综述技术主线"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    themes = [
        "多智能体任务规划与自动化工作流",
        "科学发现与自驱动实验室",
        "人机协作与人类监督",
        "垂直行业智能体(医疗、机器人)",
    ]

    y = Inches(5.6)
    for theme in themes:
        theme_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.35))
        tf = theme_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {theme}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']
        y += Inches(0.35)

def add_rq5_slide(prs):
    """RQ5：研究空白与未来挑战"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    add_header(slide, "RQ5：研究空白与未来挑战")

    # 挑战图
    img_path = os.path.join(REVIEW_FIGURES_DIR, "fig5_challenges_agenda.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.5), width=Inches(7.5))

    # 右侧详细分析
    right_x = Inches(8.2)

    # 证据来源
    evidence_title = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(4.8), Inches(0.4))
    tf = evidence_title.text_frame
    p = tf.paragraphs[0]
    p.text = "挑战识别的证据来源"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    evidence = [
        "发文量快速增长但零被引论文较多(236篇)",
        "作者合作网络较分散，过滤后仅51节点",
        "关键词中出现RAG、RL、ethics、HRI等",
        "多主题并行，尚未形成单一主导范式",
    ]

    y = Inches(2)
    for ev in evidence:
        ev_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.35))
        tf = ev_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {ev}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium']
        y += Inches(0.35)

    # 四大挑战
    challenge_title = slide.shapes.add_textbox(right_x, Inches(3.7), Inches(4.8), Inches(0.4))
    tf = challenge_title.text_frame
    p = tf.paragraphs[0]
    p.text = "四大未来挑战"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    challenges = [
        ("评价体系不统一", "不同任务、指标和评价标准导致横向比较困难，缺乏统一基准"),
        ("协作机制可解释性不足", "任务分配、冲突解决、错误溯源机制不清晰，审计困难"),
        ("垂直场景验证有限", "医疗、机器人等场景需要更强约束、审计和专家验证"),
        ("安全治理复杂化", "多智能体系统中错误传播、责任模糊、工具误用风险增加"),
    ]

    y = Inches(4.15)
    for title, desc in challenges:
        # 挑战标题 - 深黑
        title_box = slide.shapes.add_textbox(right_x, y, Inches(4.8), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {title}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(right_x + Inches(0.3), y + Inches(0.35), Inches(4.5), Inches(0.45))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['medium']

        y += Inches(0.8)

def add_conclusion_slide(prs):
    """结论总结 - 简洁深蓝风格"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 深蓝背景
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心发现总结"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 五个结论卡片
    conclusions = [
        ("爆发期已至", "2024-2025年发文量指数级增长\n2025年达344篇，同比增长325%\n领域从早期探索进入高速扩张阶段"),
        ("研究主体分散", "尚未形成稳定核心团队\n最高产作者仅发文3篇\n跨学科特征明显，全球化分布"),
        ("四层技术体系", "技术底座→协作架构→能力增强→治理应用\n研究重心转向系统协作问题\n从单模型到多智能体协同"),
        ("四条技术主线", "任务规划与自动化工作流\n科学发现与自驱动实验室\n人机协作与监督\n垂直行业智能体"),
        ("四大挑战待解", "评价体系不统一\n协作机制可解释性不足\n垂直场景验证有限\n安全治理复杂化"),
    ]

    # 第一行3个
    for i, (title, desc) in enumerate(conclusions[:3]):
        x = Inches(0.5 + i * 4.2)

        # 卡片背景 - 白色
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(3.9), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.fill.background()

        # 标题 - 深黑
        title_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(3.5), Inches(1.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['medium']

    # 第二行2个
    for i, (title, desc) in enumerate(conclusions[3:]):
        x = Inches(2.5 + i * 4.2)

        # 卡片背景 - 白色
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(3.9), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.fill.background()

        # 标题 - 深黑
        title_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(4.4), Inches(3.5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']

        # 描述 - 中灰
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(4.9), Inches(3.5), Inches(1.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['medium']

def add_acknowledgment_slide(prs):
    """致谢页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 深蓝背景
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']

    # 感谢标题
    thanks_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.5), Inches(1))
    tf = thanks_title.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听！"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # 分隔线
    sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.8), Inches(3.5), Pt(2))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = COLORS['white']
    sep_line.line.fill.background()

    # 项目信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.5), Inches(2))
    tf = info_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "GitHub: junjieg296-maker/LLM-trends-analysis-main"
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(180, 190, 210)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "数据来源：Lens.org | 有效文献：522篇"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "分析工具：Python + NetworkX + Matplotlib + Pandas"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

def add_header(slide, title):
    """添加标题栏"""
    # 顶部装饰条 - 深蓝
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 标题 - 深黑
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # 底部装饰线 - 灰色
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), SLIDE_WIDTH, Pt(2))
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLORS['border']
    bottom_line.line.fill.background()

if __name__ == "__main__":
    create_presentation()